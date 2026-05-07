"""
API Routes for Universal Business Automation Dashboard
"""
import os
import uuid
import subprocess
from datetime import datetime
from flask import Blueprint, jsonify, request, send_file, current_app
from werkzeug.utils import secure_filename

# Create blueprint
api_bp = Blueprint('api', __name__)

def allowed_file(filename):
    """Check if the file has an allowed extension"""
    return '.' in filename and \
           filename.rsplit('.', 1)[1].lower() in current_app.config['ALLOWED_EXTENSIONS']

# File type validation for better security
FILE_TYPES = {
    'pdf': ['application/pdf'],
    'docx': ['application/vnd.openxmlformats-officedocument.wordprocessingml.document'],
    'doc': ['application/msword'],
    'ppt': ['application/vnd.ms-powerpoint'],
    'pptx': ['application/vnd.openxmlformats-officedocument.presentationml.presentation'],
    'jpg': ['image/jpeg'],
    'jpeg': ['image/jpeg'],
    'png': ['image/png'],
    'bmp': ['image/bmp'],
    'tiff': ['image/tiff']
}

def validate_file_type(filename, mimetype):
    """Validate file type by both extension and MIME type"""
    if not filename:
        return False
    
    # Get file extension
    ext = filename.rsplit('.', 1)[1].lower() if '.' in filename else ''
    
    # Check if extension is allowed
    if ext not in FILE_TYPES:
        return False
    
    # If MIME type is provided, check if it matches extension
    if mimetype:
        allowed_mimes = FILE_TYPES[ext]
        return mimetype in allowed_mimes
    
    # If MIME type is empty/None, allow based on extension only (for testing)
    return True

@api_bp.route('/')
def api_root():
    """Root API endpoint with documentation"""
    return jsonify({
        'status': 'running',
        'version': '1.0.0',
        'message': 'Welcome to the Universal Business Automation API',
        'endpoints': {
            'health': '/api/health',
            'status': '/api/status',
            'leads': '/api/leads',
            'tools': '/api/tools',
            'convert': '/api/convert',
            'download': '/api/download/<file_id>'
        },
        'documentation': 'https://docs.universalbizautomat.com/api'
    })

@api_bp.route('/health')
def health():
    """Health check endpoint"""
    return jsonify({
        'status': 'healthy',
        'timestamp': datetime.utcnow().isoformat(),
        'version': '1.0.0',
        'services': {
            'database': 'connected',
            'filesystem': 'operational',
            'conversion_tools': 'available'
        }
    })

@api_bp.route('/status')
def status():
    """System status endpoint"""
    return jsonify({
        'status': 'online',
        'last_checked': datetime.utcnow().isoformat(),
        'system': {
            'upload_dir': os.path.abspath(current_app.config['UPLOAD_FOLDER']),
            'output_dir': os.path.abspath(current_app.config['OUTPUT_FOLDER']),
            'max_upload_size': f"{current_app.config['MAX_CONTENT_LENGTH'] / (1024 * 1024)}MB"
        }
    })

@api_bp.route('/leads')
def get_leads():
    """Get sample leads data"""
    # In a real app, this would come from a database
    return jsonify([
        {
            'id': 1,
            'name': 'John Doe',
            'email': 'john@example.com',
            'company': 'Acme Inc',
            'status': 'new',
            'created_at': '2025-01-10T10:30:00Z'
        },
        {
            'id': 2,
            'name': 'Jane Smith',
            'email': 'jane@example.com',
            'company': 'Globex Corp',
            'status': 'contacted',
            'created_at': '2025-01-09T14:15:00Z'
        },
        {
            'id': 3,
            'name': 'Bob Johnson',
            'email': 'bob@example.com',
            'company': 'Initech',
            'status': 'qualified',
            'created_at': '2025-01-08T09:45:00Z'
        }
    ])

@api_bp.route('/tools')
def get_tools():
    """Get available conversion tools"""
    return jsonify({
        'tools': [
            {
                'id': 'pdf_to_word',
                'name': 'PDF to Word',
                'description': 'Convert PDF documents to editable Word format',
                'input': ['pdf'],
                'output': 'docx',
                'enabled': True
            },
            {
                'id': 'word_to_pdf',
                'name': 'Word to PDF',
                'description': 'Convert Word documents to PDF format',
                'input': ['doc', 'docx'],
                'output': 'pdf',
                'enabled': True
            },
            {
                'id': 'ppt_to_pdf',
                'name': 'PowerPoint to PDF',
                'description': 'Convert PowerPoint presentations to PDF format',
                'input': ['ppt', 'pptx'],
                'output': 'pdf',
                'enabled': True
            },
            {
                'id': 'ppt_to_docx',
                'name': 'PowerPoint to Word',
                'description': 'Convert PowerPoint presentations to Word format',
                'input': ['ppt', 'pptx'],
                'output': 'docx',
                'enabled': True
            },
            {
                'id': 'image_to_pdf',
                'name': 'Image to PDF',
                'description': 'Convert image files to PDF format',
                'input': ['jpg', 'jpeg', 'png', 'bmp', 'tiff'],
                'output': 'pdf',
                'enabled': True
            },
            {
                'id': 'pdf_to_jpeg',
                'name': 'PDF to JPEG',
                'description': 'Convert PDF pages to JPEG images',
                'input': ['pdf'],
                'output': 'jpeg',
                'enabled': True
            },
            {
                'id': 'pdf_to_latex',
                'name': 'PDF to LaTeX',
                'description': 'Convert PDF documents to LaTeX source',
                'input': ['pdf'],
                'output': 'tex',
                'enabled': True
            },
            {
                'id': 'latex_to_pdf',
                'name': 'LaTeX to PDF',
                'description': 'Compile LaTeX documents to PDF',
                'input': ['tex'],
                'output': 'pdf',
                'enabled': True
            }
        ]
    })

@api_bp.route('/convert', methods=['POST'])
def convert():
    """Handle file conversion requests"""
    # Debug logging
    print(f"Request method: {request.method}")
    print(f"Request headers: {dict(request.headers)}")
    print(f"Request files: {list(request.files.keys())}")
    print(f"Request form: {dict(request.form)}")
    print(f"Request data length: {len(request.data) if request.data else 0}")
    
    # Check if the post request has the file part
    if 'file' not in request.files:
        print("ERROR: 'file' not in request.files")
        return jsonify({'error': 'No file part'}), 400
    
    file = request.files['file']
    action = request.form.get('action')
    
    # If user does not select file, browser also
    # submit an empty part without filename
    if file.filename == '':
        return jsonify({'error': 'No selected file'}), 400
    
    if not action:
        return jsonify({'error': 'No action specified'}), 400
    
    if file and allowed_file(file.filename):
        # Additional validation using MIME type
        if not validate_file_type(file.filename, file.mimetype):
            return jsonify({'error': f'Invalid file type. File {file.filename} does not match its MIME type {file.mimetype}'}), 400
        # Generate unique filename
        file_ext = file.filename.rsplit('.', 1)[1].lower()
        file_id = str(uuid.uuid4())
        filename = f"{file_id}.{file_ext}"
        
        # Save uploaded file with original extension but using file_id
        upload_path = os.path.join(current_app.config['UPLOAD_FOLDER'], filename)
        file.save(upload_path)
        
        # Prepare output path - create temp file with file_id name for conversion
        temp_filename = f"{file_id}.{file_ext}"
        temp_path = os.path.join(current_app.config['UPLOAD_FOLDER'], temp_filename)
        
        # Rename uploaded file to use file_id for consistent output naming
        os.rename(upload_path, temp_path)
        
        output_filename = f"{file_id}"
        
        # Ensure output directory exists
        if not os.path.exists(current_app.config['OUTPUT_FOLDER']):
            os.makedirs(current_app.config['OUTPUT_FOLDER'])
        
        try:
            # Perform conversion based on action
            print(f"Starting conversion: {action}")
            print(f"Input file: {temp_path}")
            print(f"Output directory: {current_app.config['OUTPUT_FOLDER']}")
            
            if action == 'pdf_to_word':
                print("Converting PDF to Word...")
                from pdf2docx import Converter
                
                output_path = os.path.join(current_app.config['OUTPUT_FOLDER'], f"{file_id}.docx")
                
                try:
                    cv = Converter(temp_path)
                    cv.convert(output_path)
                    cv.close()
                    output_filename = f"{file_id}.docx"
                    print(f"PDF to Word conversion successful: {output_filename}")
                except Exception as e:
                    print(f"PDF to Word conversion failed: {str(e)}")
                    raise Exception(f"PDF to Word conversion failed: {str(e)}")
                
            elif action == 'word_to_pdf':
                print("Converting Word to PDF...")
                # Create a more predictable output filename
                output_path = os.path.join(current_app.config['OUTPUT_FOLDER'], f"{file_id}.pdf")
                
                result = subprocess.run([
                    'soffice', '--headless', '--convert-to', 'pdf',
                    '--outdir', current_app.config['OUTPUT_FOLDER'],
                    temp_path
                ], capture_output=True, text=True)
                
                print(f"LibreOffice return code: {result.returncode}")
                print(f"LibreOffice stdout: {result.stdout}")
                print(f"LibreOffice stderr: {result.stderr}")
                
                if result.returncode != 0:
                    raise subprocess.CalledProcessError(result.returncode, result.args, result.stdout, result.stderr)
                
                # Wait for conversion and find the output file
                import time
                time.sleep(1.0)
                
                output_dir = current_app.config['OUTPUT_FOLDER']
                output_filename = None
                
                # First try to find the expected file
                if os.path.exists(output_path):
                    output_filename = f"{file_id}.pdf"
                else:
                    # Look for any PDF file created recently
                    all_files = os.listdir(output_dir)
                    pdf_files = [f for f in all_files if f.endswith('.pdf')]
                    
                    # Find the most recently created PDF
                    if pdf_files:
                        pdf_files_with_time = []
                        for f in pdf_files:
                            file_path = os.path.join(output_dir, f)
                            mtime = os.path.getmtime(file_path)
                            pdf_files_with_time.append((f, mtime))
                        
                        # Sort by modification time (newest first)
                        pdf_files_with_time.sort(key=lambda x: x[1], reverse=True)
                        output_filename = pdf_files_with_time[0][0]
                
                if not output_filename:
                    raise Exception("Could not find converted PDF file")
                
                print(f"Detected output filename: {output_filename}")
                
            elif action == 'ppt_to_docx':
                print("Converting PowerPoint to Word...")
                
                # Use python-pptx as primary method (LibreOffice CLI is unreliable on Windows)
                print("Using python-pptx for conversion...")
                
                def sanitize_text(text):
                    """Sanitize text to handle encoding issues and control characters"""
                    if not text:
                        return ""
                    
                    # Convert to string if needed
                    if not isinstance(text, str):
                        try:
                            text = str(text)
                        except:
                            return "[Unable to convert text]"
                    
                    # Remove control characters except common ones (tab, newline, carriage return)
                    import re
                    # Keep: \t (tab), \n (newline), \r (carriage return)
                    # Remove other control characters (0x00-0x1F, 0x7F-0x9F)
                    text = re.sub(r'[\x00-\x08\x0B\x0C\x0E-\x1F\x7F-\x9F]', '', text)
                    
                    # Replace any remaining problematic characters
                    text = text.replace('\x00', '')  # NULL bytes
                    text = text.replace('\uFFFD', '?')  # Replacement character
                    
                    # Strip whitespace but preserve structure
                    text = text.strip()
                    
                    return text
                
                try:
                    from pptx import Presentation
                    from docx import Document
                    
                    print(f"Loading PowerPoint file: {temp_path}")
                    
                    # Load PowerPoint presentation
                    prs = Presentation(temp_path)
                    print(f"Loaded presentation with {len(prs.slides)} slides")
                    
                    # Create Word document
                    doc = Document()
                    
                    # Extract text from each slide
                    for i, slide in enumerate(prs.slides):
                        print(f"Processing slide {i+1}")
                        
                        try:
                            # Add slide title
                            if slide.shapes.title:
                                try:
                                    title_text = sanitize_text(slide.shapes.title.text)
                                    if title_text:
                                        doc.add_heading(f"Slide {i+1}: {title_text}", level=1)
                                    else:
                                        doc.add_heading(f"Slide {i+1}", level=1)
                                except Exception as title_error:
                                    print(f"Error processing slide {i+1} title: {title_error}")
                                    doc.add_heading(f"Slide {i+1}", level=1)
                            else:
                                doc.add_heading(f"Slide {i+1}", level=1)
                            
                            # Add slide content
                            content_found = False
                            for shape in slide.shapes:
                                try:
                                    if hasattr(shape, "text") and shape != slide.shapes.title:
                                        text_content = sanitize_text(shape.text)
                                        if text_content:
                                            doc.add_paragraph(text_content)
                                            content_found = True
                                except Exception as shape_error:
                                    print(f"Error processing shape on slide {i+1}: {shape_error}")
                                    # Continue with other shapes
                                    continue
                            
                            if not content_found:
                                doc.add_paragraph("[No text content on this slide]")
                            
                            # Add page break between slides
                            if i < len(prs.slides) - 1:
                                doc.add_page_break()
                                
                        except Exception as slide_error:
                            print(f"Error processing slide {i+1}: {slide_error}")
                            # Add a placeholder for this slide and continue
                            doc.add_heading(f"Slide {i+1} - Processing Error", level=1)
                            doc.add_paragraph(f"[Error processing this slide: {str(slide_error)}]")
                            if i < len(prs.slides) - 1:
                                doc.add_page_break()
                    
                    # Save Word document
                    output_path = os.path.join(current_app.config['OUTPUT_FOLDER'], f"{file_id}.docx")
                    doc.save(output_path)
                    
                    # Verify file was created
                    if not os.path.exists(output_path):
                        raise Exception("Word document was not created")
                    
                    file_size = os.path.getsize(output_path)
                    if file_size == 0:
                        raise Exception("Word document is empty")
                    
                    output_filename = f"{file_id}.docx"
                    print(f"PowerPoint to Word conversion successful: {output_filename} ({file_size} bytes)")
                    
                except ImportError as e:
                    print(f"Required library not found: {str(e)}")
                    return jsonify({
                        'error': 'PowerPoint to Word conversion requires python-pptx and python-docx libraries',
                        'details': str(e)
                    }), 500
                except Exception as e:
                    print(f"python-pptx conversion failed: {str(e)}")
                    import traceback
                    traceback.print_exc()
                    raise Exception(f"PowerPoint to Word conversion failed: {str(e)}")
                
                                
            elif action == 'image_to_pdf':
                print("Converting Image to PDF...")
                
                # Ensure output directory exists
                output_dir = current_app.config['OUTPUT_FOLDER']
                if not os.path.exists(output_dir):
                    os.makedirs(output_dir)
                
                try:
                    from PIL import Image
                    
                    output_path = os.path.join(output_dir, f"{file_id}.pdf")
                    
                    print(f"Input image: {temp_path}")
                    print(f"Output PDF: {output_path}")
                    
                    # Open the image
                    img = Image.open(temp_path)
                    
                    # Convert to RGB if necessary (for PDF)
                    if img.mode != 'RGB':
                        img = img.convert('RGB')
                    
                    # Get image dimensions for PDF page size
                    img_width, img_height = img.size
                    
                    # Create a new PDF with the image
                    from reportlab.pdfgen import canvas
                    from reportlab.lib.pagesizes import letter
                    
                    c = canvas.Canvas(output_path, pagesize=letter)
                    page_width, page_height = letter
                    
                    # Calculate scaling to fit image on page
                    scale = min(page_width / img_width, page_height / img_height) * 0.9  # 90% of page
                    scaled_width = img_width * scale
                    scaled_height = img_height * scale
                    
                    # Center the image on the page
                    x = (page_width - scaled_width) / 2
                    y = (page_height - scaled_height) / 2
                    
                    # Draw the image
                    c.drawImage(temp_path, x, y, scaled_width, scaled_height)
                    c.save()
                    
                    output_filename = f"{file_id}.pdf"
                    print(f"Image to PDF conversion successful: {output_filename}")
                    
                except ImportError:
                    print("ReportLab not available, falling back to PIL PDF save...")
                    try:
                        from PIL import Image
                        
                        # Open the image
                        img = Image.open(temp_path)
                        
                        # Convert to RGB if necessary (for PDF)
                        if img.mode != 'RGB':
                            img = img.convert('RGB')
                        
                        # Create PDF output path
                        output_path = os.path.join(output_dir, f"{file_id}.pdf")
                        
                        # Save as PDF
                        img.save(output_path, "PDF", resolution=100.0)
                        
                        output_filename = f"{file_id}.pdf"
                        print(f"Image to PDF conversion successful with PIL: {output_filename}")
                        
                    except ImportError:
                        raise Exception("Neither ReportLab nor PIL available for image to PDF conversion")
                except Exception as e:
                    print(f"ReportLab conversion failed: {str(e)}")
                    # Fallback to PIL method
                    try:
                        from PIL import Image
                        
                        # Open the image
                        img = Image.open(temp_path)
                        
                        # Convert to RGB if necessary (for PDF)
                        if img.mode != 'RGB':
                            img = img.convert('RGB')
                        
                        # Create PDF output path
                        output_path = os.path.join(output_dir, f"{file_id}.pdf")
                        
                        # Save as PDF
                        img.save(output_path, "PDF", resolution=100.0)
                        
                        output_filename = f"{file_id}.pdf"
                        print(f"Image to PDF conversion successful with PIL fallback: {output_filename}")
                        
                    except ImportError:
                        raise Exception("No image processing libraries available for image to PDF conversion")
                
            elif action == 'pdf_to_jpeg':
                subprocess.run([
                    'convert', '-density', '300',
                    temp_path,
                    os.path.join(current_app.config['OUTPUT_FOLDER'], f"{file_id}.jpg")
                ], check=True)
                output_filename += '.jpg'
                
            elif action == 'pdf_to_latex':
                subprocess.run([
                    'pandoc', temp_path, '-o', os.path.join(current_app.config['OUTPUT_FOLDER'], f"{file_id}.tex")
                ], check=True)
                output_filename += '.tex'
                
            elif action == 'latex_to_pdf':
                subprocess.run([
                    'pdflatex', '-output-directory', current_app.config['OUTPUT_FOLDER'],
                    temp_path
                ], check=True)
                output_filename += '.pdf'
                
            else:
                return jsonify({'error': 'Invalid action'}), 400
                
            # Clean up temp file
            try:
                os.remove(temp_path)
            except:
                pass
                
            # Debug logging
            print(f"Conversion successful. Output filename: {output_filename}")
            expected_path = os.path.join(current_app.config['OUTPUT_FOLDER'], output_filename)
            print(f"Expected file path: {expected_path}")
            print(f"File exists: {os.path.exists(expected_path)}")
            
            return jsonify({
                'success': True,
                'message': 'Conversion successful',
                'file_id': file_id,
                'download_url': f"/api/download/{output_filename}"
            })
            
        except subprocess.CalledProcessError as e:
            # Clean up temp file on error
            try:
                os.remove(temp_path)
            except:
                pass
            current_app.logger.error(f"Conversion failed: {str(e)}")
            return jsonify({
                'error': 'Conversion failed',
                'details': str(e)
            }), 500
            
        except Exception as e:
            current_app.logger.error(f"Unexpected error: {str(e)}")
            return jsonify({
                'error': 'An unexpected error occurred',
                'details': str(e)
            }), 500
            
    return jsonify({'error': 'File type not allowed'}), 400

@api_bp.route('/download/<path:filename>')
def download_file(filename):
    """Download converted files"""
    # Debug logging
    print(f"Download request for filename: {filename}")
    
    # Sanitize filename to prevent directory traversal
    safe_filename = secure_filename(filename)
    file_path = os.path.join(current_app.config['OUTPUT_FOLDER'], safe_filename)
    
    print(f"Safe filename: {safe_filename}")
    print(f"Full file path: {file_path}")
    print(f"Output folder: {current_app.config['OUTPUT_FOLDER']}")
    print(f"Files in output folder: {os.listdir(current_app.config['OUTPUT_FOLDER'])}")
    print(f"File exists: {os.path.exists(file_path)}")
    
    if not os.path.exists(file_path):
        return jsonify({'error': 'File not found'}), 404
        
    return send_file(
        file_path,
        as_attachment=True,
        download_name=safe_filename
    )
